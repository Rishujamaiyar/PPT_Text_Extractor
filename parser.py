!pip install python-docx
!pip install python-pptx
!pip install pysbd

import os
import sys
import traceback

from pptx import Presentation
import pysbd
from pathlib import Path
import re
from openpyxl import load_workbook
from nltk.corpus import stopwords
import docx
import json
import logging
from pptx import Presentation
from pptx.chart.chart import Chart
from pptx.enum.shapes import PP_PLACEHOLDER
from datetime import datetime
from configparser import ConfigParser



def extract_from_ppt(file):
    prs = Presentation(file)
    i = 1
    
    file_extract = {}
    text_extract = []
    title_list = []
    footer_list = []
    
    for slide in prs.slides:
        
        #EXTRACT TITLE
        temp_dict = dict()
        title = ""
        if slide.shapes:
            if slide.shapes.title:
                title = slide.shapes.title.text
        temp_dict["slide_%s_title" % i] = title
        title_list.append(temp_dict)
        
        #EXTRACT FOOTER
        temp_dict = dict()
        text_value = ""
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.type == PP_PLACEHOLDER.FOOTER:
                if placeholder:
                    text_value = placeholder.text
            temp_dict["slide_%s_footer" % i] = text_value
        footer_list.append(temp_dict)
        
        #EXTRACT BODY
        temp_dict = dict()
        text_runs = list()
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                paragraph_list = list()
                for run in paragraph.runs:
                    paragraph_list.append(run.text)
                if paragraph_list:
                    paragraph_text = " ".join(paragraph_list)
                    text_runs.append(paragraph_text)
                        
        #EXTRACT BODY TABLE CONTENTS              
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            tbl = shape.table
            row_count = len(tbl.rows)
            col_count = len(tbl.columns)
            for r in range(0, row_count):
                for c in range(0, col_count):
                    cell = tbl.cell(r, c)
                    paragraphs = cell.text_frame.paragraphs
                    full_text_list = list()
                    for paragraph in paragraphs:
                        for run in paragraph.runs:
                            full_text_list.append(run.text)
                    if full_text_list:
                        full_text = " ".join(full_text_list)
                        text_runs.append(full_text)
                            
            cleaned_test_runs = list()
            for the_extract in text_runs:
                for footer_dtl in footer_list:
                    if the_extract not in list(footer_dtl.values()):
                        cleaned_test_runs.append(the_extract)
                        
            temp_dict["slide_%s_body" % i] = list(set(cleaned_test_runs))
            text_extract.append(temp_dict)
        i += 1
            
    file_extract["title"] = title_list
    file_extract["footer"] = footer_list
    file_extract["body"] = text_extract 
    return file_extract
